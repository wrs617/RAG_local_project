import time
from io import BytesIO
import zipfile
import xml.etree.ElementTree as ET

import streamlit as st
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from knowledge_base import KnowledgeBaseService


def decode_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "gbk", "gb18030"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_docx_text(file_bytes: bytes) -> str:
    # 使用 docx 的 XML 结构提取文本，避免额外安装 python-docx 依赖
    with zipfile.ZipFile(BytesIO(file_bytes)) as zf:
        xml_data = zf.read("word/document.xml")
    root = ET.fromstring(xml_data)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    texts = [node.text for node in root.findall(".//w:t", ns) if node.text]
    return "\n".join(texts)


def extract_pdf_text(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    parts = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if txt.strip():
            parts.append(txt)
    return "\n".join(parts)


def extract_ppt_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


def extract_xlsx_text(file_bytes: bytes) -> str:
    wb = load_workbook(filename=BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet] {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                parts.append(" | ".join(values))
    wb.close()
    return "\n".join(parts)


def extract_text_by_extension(file_name: str, file_bytes: bytes) -> str:
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    if ext == "txt":
        return decode_text(file_bytes)
    if ext == "docx":
        return extract_docx_text(file_bytes)
    if ext == "pdf":
        return extract_pdf_text(file_bytes)
    if ext in ("ppt", "pptx"):
        return extract_ppt_text(file_bytes)
    if ext == "xlsx":
        return extract_xlsx_text(file_bytes)

    raise ValueError(f"不支持的文件格式: {ext}")


st.title("知识库更新服务")

uploader_file = st.file_uploader(
    "请上传文件（txt/docx/pdf/ppt/pptx/xlsx）",
    type=["txt", "docx", "pdf", "ppt", "pptx", "xlsx"],
    accept_multiple_files=False,
)

if "service" not in st.session_state:
    st.session_state["service"] = KnowledgeBaseService()

if uploader_file is not None:
    file_name = uploader_file.name
    file_type = uploader_file.type
    file_size = uploader_file.size / 1024

    st.subheader(f"文件名: {file_name}")
    st.write(f"文件格式: {file_type}")
    st.write(f"文件大小: {file_size:.2f} KB")

    if st.button("上传到知识库", type="primary"):
        try:
            raw_bytes = uploader_file.getvalue()
            text = extract_text_by_extension(file_name, raw_bytes)
            if not text.strip():
                st.warning("提取到的内容为空，无法入库。")
            else:
                with st.spinner("载入知识库中..."):
                    time.sleep(0.3)
                    result = st.session_state["service"].upload_by_str(text, file_name)
                st.success(result)
        except Exception as err:
            st.error(f"上传失败: {type(err).__name__}: {err}")