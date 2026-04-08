import uuid
from io import BytesIO
import zipfile
import xml.etree.ElementTree as ET

import streamlit as st

from knowledge_base import KnowledgeBaseService
from rag import RagService


def decode_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "gbk", "gb18030"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_docx_text(file_bytes: bytes) -> str:
    with zipfile.ZipFile(BytesIO(file_bytes)) as zf:
        xml_data = zf.read("word/document.xml")
    root = ET.fromstring(xml_data)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    texts = [node.text for node in root.findall(".//w:t", ns) if node.text]
    return "\n".join(texts)


def extract_pdf_text(file_bytes: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(file_bytes))
    parts = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if txt.strip():
            parts.append(txt)
    return "\n".join(parts)


def extract_ppt_text(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


def extract_xlsx_text(file_bytes: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet] {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                parts.append(" | ".join(values))
    wb.close()
    return "\n".join(parts)


def extract_text_by_extension(file_name: str, file_bytes: bytes) -> str:
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    if ext == "txt":
        return decode_text(file_bytes)
    if ext == "docx":
        return extract_docx_text(file_bytes)
    if ext == "pdf":
        return extract_pdf_text(file_bytes)
    if ext in ("ppt", "pptx"):
        return extract_ppt_text(file_bytes)
    if ext == "xlsx":
        return extract_xlsx_text(file_bytes)

    raise ValueError(f"不支持的文件格式: {ext}")


def format_error(err: Exception) -> str:
    return f"{type(err).__name__}: {err}"


st.set_page_config(page_title="本地RAG", layout="wide")
st.title("本地RAG")
st.divider()

if "message" not in st.session_state:
    st.session_state["message"] = [
        {"role": "assistant", "content": "你好，有什么可以帮助你？"}
    ]

if "session_id" not in st.session_state:
    st.session_state["session_id"] = str(uuid.uuid4())

if "rag" not in st.session_state:
    st.session_state["rag"] = RagService()

if "kb_service" not in st.session_state:
    st.session_state["kb_service"] = KnowledgeBaseService()

with st.sidebar:
    st.subheader("文件上传")
    uploader_file = st.file_uploader(
        "请上传文件（txt/docx/pdf/ppt/pptx/xlsx）",
        type=["txt", "docx", "pdf", "ppt", "pptx", "xlsx"],
        accept_multiple_files=False,
    )

    if uploader_file is not None:
        st.caption(f"文件名: {uploader_file.name}")
        st.caption(f"文件大小: {uploader_file.size / 1024:.2f} KB")

        if st.button("上传到知识库", type="primary"):
            try:
                text = extract_text_by_extension(uploader_file.name, uploader_file.getvalue())
                if not text.strip():
                    st.warning("提取到的内容为空，无法入库。")
                else:
                    with st.spinner("正在写入知识库..."):
                        result = st.session_state["kb_service"].upload_by_str(text, uploader_file.name)
            except Exception as err:
                st.error(f"上传失败：{format_error(err)}")
            else:
                if "成功" in result:
                    st.session_state["rag"] = RagService()
                    st.success(result)
                elif "跳过" in result:
                    st.info(result)
                else:
                    st.warning(result)

    st.divider()
    if st.button("新建会话"):
        st.session_state["session_id"] = str(uuid.uuid4())
        st.session_state["message"] = [
            {"role": "assistant", "content": "你好，有什么可以帮助你？"}
        ]
        st.rerun()

for message in st.session_state["message"]:
    st.chat_message(message["role"]).write(message["content"])

prompt = st.chat_input("请输入你的问题")

if prompt:
    st.chat_message("user").write(prompt)
    st.session_state["message"].append({"role": "user", "content": prompt})

    ai_res_list = []
    session_config = {"configurable": {"session_id": st.session_state["session_id"]}}

    try:
        with st.spinner("AI思考中..."):
            res_stream = st.session_state["rag"].chain.stream(
                {"input": prompt},
                config=session_config,
            )

            def capture(generator, cache_list):
                for chunk in generator:
                    cache_list.append(chunk)
                    yield chunk

            st.chat_message("assistant").write_stream(capture(res_stream, ai_res_list))
    except Exception as err:
        err_msg = f"问答失败：{format_error(err)}"
        st.chat_message("assistant").write(err_msg)
        st.session_state["message"].append({"role": "assistant", "content": err_msg})
    else:
        final_answer = "".join(ai_res_list).strip() or "当前没有生成有效回复，请重试一次。"
        st.session_state["message"].append({"role": "assistant", "content": final_answer})